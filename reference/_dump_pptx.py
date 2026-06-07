# -*- coding: utf-8 -*-
import sys
from pptx import Presentation

OUT = r'c:\myProjects\logisticslink\reference\_pptx_dump.txt'
SRC = r'c:\myProjects\logisticslink\reference\20240322_DFP_Process_v0.39.pptx'

p = Presentation(SRC)
lines = []
lines.append(f'Total slides: {len(p.slides)}')
lines.append('=' * 80)

for i, slide in enumerate(p.slides, 1):
    lines.append(f'\n--- Slide {i} ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = ''.join(run.text for run in para.runs)
                if txt.strip():
                    lines.append(txt)
        elif shape.has_table:
            try:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    lines.append(' | '.join(cells))
            except Exception as e:
                lines.append(f'[table error: {e}]')

with open(OUT, 'w', encoding='utf-8') as f:
    f.write('\n'.join(lines))

print(f'Wrote {len(lines)} lines to {OUT}')
